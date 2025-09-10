"""
文档解析器
支持多种文档格式的解析和图片识别
"""

import logging
import os
import time
import uuid
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path

# PDF解析
import PyPDF2
import pdfplumber
import fitz  # PyMuPDF

# Office文档解析
from pptx import Presentation
import docx

# 图片处理
from PIL import Image

from .document_models import DocumentType, DocumentInfo, ParseResult, ImageAnalysisResult

logger = logging.getLogger(__name__)


class DocumentParser:
    """文档解析器类"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf': self._parse_pdf,
            'pptx': self._parse_pptx,
            'ppt': self._parse_ppt,
            'docx': self._parse_docx,
            'doc': self._parse_doc,
            'txt': self._parse_txt,
            'png': self._parse_image,
            'jpg': self._parse_image,
            'jpeg': self._parse_image,
            'gif': self._parse_image,
            'bmp': self._parse_image
        }
    
    def parse_document(self, file_path: str, document_info: DocumentInfo) -> ParseResult:
        """
        解析文档
        Args:
            file_path: 文件路径
            document_info: 文档信息
        Returns:
            ParseResult: 解析结果
        """
        start_time = time.time()
        
        try:
            file_ext = document_info.document_type.value.lower()
            
            if file_ext not in self.supported_formats:
                return ParseResult(
                    document_id=document_info.id,
                    success=False,
                    error_message=f"不支持的文件格式: {file_ext}"
                )
            
            # 调用对应的解析方法
            parser_func = self.supported_formats[file_ext]
            result = parser_func(file_path, document_info)
            
            # 计算解析耗时
            duration = time.time() - start_time
            result.duration = duration
            
            # 计算字数
            if result.text_content:
                result.word_count = len(result.text_content.replace(' ', ''))
            
            logger.info(f"文档解析完成: {document_info.filename}, 耗时: {duration:.2f}s, 字数: {result.word_count}")
            
            return result
            
        except Exception as e:
            logger.error(f"文档解析失败: {document_info.filename}, 错误: {e}")
            return ParseResult(
                document_id=document_info.id,
                success=False,
                error_message=str(e),
                duration=time.time() - start_time
            )
    
    def _parse_pdf(self, file_path: str, document_info: DocumentInfo) -> ParseResult:
        """解析PDF文件"""
        text_content = ""
        page_count = 0
        images = []
        tables = []
        
        try:
            # 优先使用pdfplumber
            text_content, page_count, tables = self._parse_pdf_with_pdfplumber(file_path)
            
            # 如果pdfplumber解析失败，使用PyMuPDF
            if not text_content.strip():
                text_content, page_count, images = self._parse_pdf_with_pymupdf(file_path)
            
            # 如果仍然失败，使用PyPDF2
            if not text_content.strip():
                text_content, page_count = self._parse_pdf_with_pypdf2(file_path)
            
            return ParseResult(
                document_id=document_info.id,
                success=True,
                text_content=text_content,
                page_count=page_count,
                images=images,
                tables=tables,
                metadata={
                    "parser": "multi",
                    "file_size": document_info.file_size
                }
            )
            
        except Exception as e:
            logger.error(f"PDF解析失败: {e}")
            return ParseResult(
                document_id=document_info.id,
                success=False,
                error_message=str(e)
            )
    
    def _parse_pdf_with_pdfplumber(self, file_path: str) -> Tuple[str, int, List[Dict]]:
        """使用pdfplumber解析PDF"""
        text_content = ""
        page_count = 0
        tables = []
        
        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)
            
            for page in pdf.pages:
                # 提取文本
                page_text = page.extract_text()
                if page_text:
                    text_content += page_text + "\n\n"
                
                # 提取表格
                page_tables = page.extract_tables()
                if page_tables:
                    for table in page_tables:
                        tables.append({
                            "page": page.page_number,
                            "data": table
                        })
        
        return text_content.strip(), page_count, tables
    
    def _parse_pdf_with_pymupdf(self, file_path: str) -> Tuple[str, int, List[str]]:
        """使用PyMuPDF解析PDF"""
        text_content = ""
        page_count = 0
        images = []
        
        doc = fitz.open(file_path)
        page_count = doc.page_count
        
        for page_num in range(page_count):
            page = doc[page_num]
            
            # 提取文本
            text_content += page.get_text() + "\n\n"
            
            # 提取图片
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    
                    if pix.n - pix.alpha < 4:  # 确保是RGB格式
                        img_path = f"extracted_images/{uuid.uuid4()}_{page_num}_{img_index}.png"
                        os.makedirs(os.path.dirname(img_path), exist_ok=True)
                        pix.save(img_path)
                        images.append(img_path)
                    
                    pix = None
                except Exception as e:
                    logger.warning(f"提取图片失败: {e}")
        
        doc.close()
        return text_content.strip(), page_count, images
    
    def _parse_pdf_with_pypdf2(self, file_path: str) -> Tuple[str, int]:
        """使用PyPDF2解析PDF"""
        text_content = ""
        page_count = 0
        
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            page_count = len(reader.pages)
            
            for page in reader.pages:
                text_content += page.extract_text() + "\n\n"
        
        return text_content.strip(), page_count
    
    def _parse_pptx(self, file_path: str, document_info: DocumentInfo) -> ParseResult:
        """解析PPTX文件"""
        try:
            prs = Presentation(file_path)
            text_content = ""
            page_count = len(prs.slides)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"幻灯片 {slide_num}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                text_content += slide_text + "\n\n"
            
            return ParseResult(
                document_id=document_info.id,
                success=True,
                text_content=text_content.strip(),
                page_count=page_count,
                metadata={
                    "parser": "python-pptx",
                    "slide_count": page_count
                }
            )
            
        except Exception as e:
            logger.error(f"PPTX解析失败: {e}")
            return ParseResult(
                document_id=document_info.id,
                success=False,
                error_message=str(e)
            )
    
    def _parse_ppt(self, file_path: str, document_info: DocumentInfo) -> ParseResult:
        """解析PPT文件（旧格式）"""
        # PPT格式较复杂，这里提供基础实现
        # 实际项目中可能需要使用win32com（Windows）或LibreOffice
        logger.warning("PPT格式解析暂不完全支持，建议转换为PPTX格式")
        
        return ParseResult(
            document_id=document_info.id,
            success=False,
            error_message="PPT格式解析暂不支持，请转换为PPTX格式"
        )
    
    def _parse_docx(self, file_path: str, document_info: DocumentInfo) -> ParseResult:
        """解析DOCX文件"""
        try:
            doc = docx.Document(file_path)
            text_content = ""
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content += paragraph.text + "\n"
            
            # 解析表格
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append({"data": table_data})
            
            return ParseResult(
                document_id=document_info.id,
                success=True,
                text_content=text_content.strip(),
                page_count=1,  # DOCX没有固定页数概念
                tables=tables,
                metadata={
                    "parser": "python-docx",
                    "paragraph_count": len(doc.paragraphs),
                    "table_count": len(doc.tables)
                }
            )
            
        except Exception as e:
            logger.error(f"DOCX解析失败: {e}")
            return ParseResult(
                document_id=document_info.id,
                success=False,
                error_message=str(e)
            )
    
    def _parse_doc(self, file_path: str, document_info: DocumentInfo) -> ParseResult:
        """解析DOC文件（旧格式）"""
        logger.warning("DOC格式解析暂不完全支持，建议转换为DOCX格式")
        
        return ParseResult(
            document_id=document_info.id,
            success=False,
            error_message="DOC格式解析暂不支持，请转换为DOCX格式"
        )
    
    def _parse_txt(self, file_path: str, document_info: DocumentInfo) -> ParseResult:
        """解析TXT文件"""
        try:
            # 尝试多种编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            text_content = ""
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text_content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if not text_content:
                raise ValueError("无法解码文本文件")
            
            return ParseResult(
                document_id=document_info.id,
                success=True,
                text_content=text_content,
                page_count=1,
                metadata={
                    "parser": "text",
                    "encoding": encoding
                }
            )
            
        except Exception as e:
            logger.error(f"TXT解析失败: {e}")
            return ParseResult(
                document_id=document_info.id,
                success=False,
                error_message=str(e)
            )
    
    def _parse_image(self, file_path: str, document_info: DocumentInfo) -> ParseResult:
        """解析图片文件"""
        try:
            # 基础图片信息
            with Image.open(file_path) as img:
                width, height = img.size
                format_name = img.format
            
            # 这里可以集成OCR服务进行文字识别
            # 暂时返回基础信息
            metadata = {
                "parser": "PIL",
                "width": width,
                "height": height,
                "format": format_name
            }
            
            return ParseResult(
                document_id=document_info.id,
                success=True,
                text_content="",  # OCR结果将在这里
                page_count=1,
                images=[file_path],
                metadata=metadata
            )
            
        except Exception as e:
            logger.error(f"图片解析失败: {e}")
            return ParseResult(
                document_id=document_info.id,
                success=False,
                error_message=str(e)
            )
    
    def analyze_image_with_llm(self, image_path: str, llm_client=None) -> ImageAnalysisResult:
        """
        使用大模型分析图片内容
        Args:
            image_path: 图片路径
            llm_client: LLM客户端（支持视觉的模型，如GPT-4V、豆包等）
        Returns:
            ImageAnalysisResult: 分析结果
        """
        try:
            if not llm_client:
                return ImageAnalysisResult(
                    image_path=image_path,
                    text_content="",
                    confidence=0.0
                )
            
            # 这里集成视觉大模型进行图片分析
            # 具体实现依赖于所使用的模型
            
            # 示例：调用豆包视觉模型
            # result = llm_client.analyze_image(image_path)
            
            return ImageAnalysisResult(
                image_path=image_path,
                text_content="",  # 模型分析结果
                confidence=0.0
            )
            
        except Exception as e:
            logger.error(f"图片分析失败: {e}")
            return ImageAnalysisResult(
                image_path=image_path,
                text_content="",
                confidence=0.0
            )
    
    def get_supported_formats(self) -> List[str]:
        """获取支持的文件格式"""
        return list(self.supported_formats.keys())
    
    def is_supported_format(self, file_extension: str) -> bool:
        """检查是否支持指定格式"""
        return file_extension.lower() in self.supported_formats
