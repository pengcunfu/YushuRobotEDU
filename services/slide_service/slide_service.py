"""
幻灯片处理核心服务
支持PPT、PPTX、PDF等格式的解析和处理
"""

import os
import uuid
import base64
import logging
from io import BytesIO
from pathlib import Path
from typing import List, Optional, Dict, Any, Tuple
import asyncio

# 导入第三方库
try:
    from PIL import Image
    import fitz  # PyMuPDF
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    
try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

from .slide_models import SlideInfo, SlideContent, SlideNarration, SlideProcessRequest
from ..llm_service.llm_manager import llm_manager
from ..llm_service.chat_message import ChatMessage
from ..tts_service.tts_manager import tts_manager
from .ppt_converter import ppt_converter

try:
    from ..config_service.config_center_service import config_center_service
except ImportError:
    config_center_service = None

logger = logging.getLogger(__name__)

class SlideService:
    """幻灯片处理服务"""
    
    def __init__(self, storage_path: str = "uploads/slides"):
        self.storage_path = Path(storage_path)
        self.storage_path.mkdir(parents=True, exist_ok=True)
        
        # 图片存储路径
        self.images_path = self.storage_path / "images"
        self.images_path.mkdir(exist_ok=True)
        
        # 缩略图存储路径
        self.thumbnails_path = self.storage_path / "thumbnails"
        self.thumbnails_path.mkdir(exist_ok=True)
        
        # 音频存储路径
        self.audio_path = self.storage_path / "audio"
        self.audio_path.mkdir(exist_ok=True)
    
    async def process_document(self, request: SlideProcessRequest) -> SlideContent:
        """处理文档，提取幻灯片内容"""
        try:
            file_path = Path(request.file_path)
            if not file_path.exists():
                raise FileNotFoundError(f"文件不存在: {request.file_path}")
            
            file_extension = file_path.suffix.lower()
            
            if file_extension == '.pdf':
                return await self._process_pdf(request)
            elif file_extension in ['.pptx', '.ppt']:
                return await self._process_pptx(request)
            else:
                raise ValueError(f"不支持的文件格式: {file_extension}")
                
        except Exception as e:
            logger.error(f"文档处理失败: {str(e)}")
            raise
    
    async def _process_pdf(self, request: SlideProcessRequest) -> SlideContent:
        """处理PDF文件"""
        if not PDF_SUPPORT:
            raise RuntimeError("PDF处理库未安装，请安装 pdf2image 和 PyMuPDF")
        
        file_path = Path(request.file_path)
        filename = file_path.name
        
        # 使用PyMuPDF打开PDF
        doc = fitz.open(str(file_path))
        total_pages = len(doc)
        
        slides = []
        
        for page_num in range(total_pages):
            page = doc.load_page(page_num)
            
            # 提取文字内容
            text_content = page.get_text()
            
            # 生成幻灯片图片
            image_url = None
            thumbnail_url = None
            
            if request.extract_images:
                # 将页面转换为图片
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2倍分辨率
                img_data = pix.tobytes("png")
                
                # 保存原图
                image_filename = f"{request.document_id}_page_{page_num + 1}.png"
                image_path = self.images_path / image_filename
                
                with open(image_path, "wb") as f:
                    f.write(img_data)
                
                image_url = f"/api/slides/images/{image_filename}"
                
                # 生成缩略图
                if request.generate_thumbnails:
                    thumbnail_url = await self._generate_thumbnail(image_path, page_num + 1, request.document_id)
            
            # 尝试提取标题（第一行文字）
            title = None
            if text_content.strip():
                lines = text_content.strip().split('\n')
                title = lines[0][:50] if lines else f"第{page_num + 1}页"
            else:
                title = f"第{page_num + 1}页"
            
            slide_info = SlideInfo(
                id=f"{request.document_id}_slide_{page_num + 1}",
                page_number=page_num + 1,
                title=title,
                content=text_content,
                image_url=image_url,
                thumbnail_url=thumbnail_url
            )
            
            slides.append(slide_info)
        
        doc.close()
        
        return SlideContent(
            document_id=request.document_id,
            filename=filename,
            total_pages=total_pages,
            slides=slides
        )
    
    async def _process_pptx(self, request: SlideProcessRequest) -> SlideContent:
        """处理PPTX文件"""
        file_path = Path(request.file_path)
        filename = file_path.name
        
        # 首先尝试使用PPT转图片服务
        slides = []
        total_pages = 0
        
        try:
            # 使用外部PPT转图片服务
            logger.info(f"使用PPT转图片服务处理文件: {file_path}")
            
            image_paths = await ppt_converter.convert_ppt_to_images(
                ppt_path=str(file_path),
                document_id=request.document_id,
                output_dir=str(self.images_path)
            )
            
            total_pages = len(image_paths)
            logger.info(f"PPT转图片服务成功生成 {total_pages} 张图片")
            
            # 如果有python-pptx支持，尝试提取文本内容
            text_contents = []
            titles = []
            notes_list = []
            
            if PPTX_SUPPORT:
                try:
                    prs = Presentation(str(file_path))
                    for slide_num, slide in enumerate(prs.slides):
                        # 提取标题
                        title = None
                        if slide.shapes.title:
                            title = slide.shapes.title.text
                        titles.append(title or f"第{slide_num + 1}页")
                        
                        # 提取文字内容
                        text_content = ""
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                text_content += shape.text + "\n"
                        text_contents.append(text_content.strip())
                        
                        # 提取演讲者备注
                        notes = ""
                        if slide.has_notes_slide:
                            notes_slide = slide.notes_slide
                            for shape in notes_slide.shapes:
                                if hasattr(shape, "text") and shape.text:
                                    notes += shape.text + "\n"
                        notes_list.append(notes.strip() if notes else None)
                except Exception as e:
                    logger.warning(f"提取PPTX文本内容失败: {e}")
                    # 创建默认文本内容
                    for i in range(total_pages):
                        titles.append(f"第{i + 1}页")
                        text_contents.append("")
                        notes_list.append(None)
            else:
                # 没有python-pptx支持，创建默认文本内容
                for i in range(total_pages):
                    titles.append(f"第{i + 1}页")
                    text_contents.append("")
                    notes_list.append(None)
            
            # 创建幻灯片信息
            for i, image_path in enumerate(image_paths):
                image_filename = Path(image_path).name
                image_url = f"/api/slides/images/{image_filename}" if request.extract_images else None
                
                # 生成缩略图
                thumbnail_url = None
                if request.generate_thumbnails and image_url:
                    thumbnail_url = await self._generate_thumbnail(
                        Path(image_path), i + 1, request.document_id
                    )
                
                slide_info = SlideInfo(
                    id=f"{request.document_id}_slide_{i + 1}",
                    page_number=i + 1,
                    title=titles[i] if i < len(titles) else f"第{i + 1}页",
                    content=text_contents[i] if i < len(text_contents) else "",
                    image_url=image_url,
                    thumbnail_url=thumbnail_url,
                    notes=notes_list[i] if i < len(notes_list) else None
                )
                
                slides.append(slide_info)
                
        except Exception as e:
            logger.error(f"PPT转图片服务失败: {e}")
            
            # 如果PPT转图片服务失败，只提取文本内容，不生成图片
            if not PPTX_SUPPORT:
                raise RuntimeError("PPT转图片服务不可用，且PPTX处理库未安装，请安装 python-pptx")
            
            # 只提取文本内容，不生成图片
            logger.info("PPT转图片服务不可用，仅提取文本内容")
            prs = Presentation(str(file_path))
            total_pages = len(prs.slides)
            
            for slide_num, slide in enumerate(prs.slides):
                # 提取文字内容
                text_content = ""
                title = None
                
                # 提取标题
                if slide.shapes.title:
                    title = slide.shapes.title.text
                
                # 提取所有文字内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content += shape.text + "\n"
                
                # 提取演讲者备注
                notes = ""
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    for shape in notes_slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            notes += shape.text + "\n"
                
                if not title:
                    title = f"第{slide_num + 1}页"
                
                slide_info = SlideInfo(
                    id=f"{request.document_id}_slide_{slide_num + 1}",
                    page_number=slide_num + 1,
                    title=title,
                    content=text_content.strip(),
                    image_url=None,  # 不生成图片
                    thumbnail_url=None,
                    notes=notes.strip() if notes else None
                )
                
                slides.append(slide_info)
        
        return SlideContent(
            document_id=request.document_id,
            filename=filename,
            total_pages=total_pages,
            slides=slides
        )
    
    
    async def _generate_thumbnail(self, image_path: Path, page_num: int, document_id: str) -> str:
        """生成缩略图"""
        try:
            # 打开原图
            with Image.open(image_path) as img:
                # 生成缩略图 (150x100)
                img.thumbnail((150, 100), Image.Resampling.LANCZOS)
                
                # 保存缩略图
                thumbnail_filename = f"{document_id}_thumb_{page_num}.png"
                thumbnail_path = self.thumbnails_path / thumbnail_filename
                
                img.save(thumbnail_path, "PNG")
                
                return f"/api/slides/thumbnails/{thumbnail_filename}"
                
        except Exception as e:
            logger.error(f"缩略图生成失败: {str(e)}")
            return None
    
    async def generate_narration(self, slide_info: SlideInfo, style: str = "professional") -> SlideNarration:
        """为幻灯片生成AI讲解内容"""
        try:
            # 构建讲解提示词
            prompt = self._build_narration_prompt(slide_info, style)
            
            # 尝试调用LLM服务生成讲解内容
            try:
                narration = await self._generate_llm_narration(prompt)
            except Exception as llm_error:
                logger.warning(f"LLM服务不可用，使用基础讲解: {str(llm_error)}")
                narration = await self._generate_basic_narration(slide_info, style)
            
            return SlideNarration(
                slide_id=slide_info.id,
                narration=narration,
                duration=len(narration) * 0.15  # 估算阅读时长
            )
            
        except Exception as e:
            logger.error(f"讲解生成失败: {str(e)}")
            raise
    
    def _build_narration_prompt(self, slide_info: SlideInfo, style: str) -> str:
        """构建讲解提示词"""
        style_prompts = {
            "professional": "请用专业、正式的语调",
            "casual": "请用轻松、友好的语调", 
            "educational": "请用教育性、解释性的语调",
            "enthusiastic": "请用热情、激励性的语调"
        }
        
        style_instruction = style_prompts.get(style, style_prompts["professional"])
        
        prompt = f"""
{style_instruction}为以下幻灯片内容生成讲解文稿：

幻灯片标题：{slide_info.title or '无标题'}
幻灯片内容：
{slide_info.content or '无内容'}

讲解要求：
1. 语言自然流畅，适合口语表达
2. 突出重点内容
3. 时长控制在30-60秒
4. 避免直接重复幻灯片文字
5. 可以适当扩展和解释内容

请生成讲解文稿：
"""
        return prompt
    
    async def _generate_llm_narration(self, prompt: str) -> str:
        """调用LLM服务生成讲解内容"""
        try:
            # 获取可用的LLM提供商
            available_providers = llm_manager.get_available_providers()
            if not available_providers:
                raise Exception("没有可用的LLM提供商")
            
            # 使用第一个可用的提供商
            provider = available_providers[0]
            logger.info(f"使用 {provider} LLM生成讲解内容")
            
            # 调用LLM生成讲解
            result = llm_manager.chat(
                provider=provider,
                message=prompt,
                temperature=0.7,
                max_tokens=500
            )
            
            if not result.get('success', False):
                raise Exception(f"LLM调用失败: {result.get('error', '未知错误')}")
            
            # 提取生成的文本
            response = result.get('response', '')
            if isinstance(response, dict):
                # 如果response是字典，尝试提取文本内容
                narration = response.get('content', '') or response.get('text', '') or str(response)
            else:
                narration = str(response)
            
            if not narration.strip():
                raise Exception("LLM返回空内容")
            
            return narration.strip()
            
        except Exception as e:
            logger.error(f"LLM服务调用失败: {str(e)}")
            raise
    
    async def _generate_basic_narration(self, slide_info: SlideInfo, style: str) -> str:
        """生成模拟讲解内容（实际应用中应调用LLM服务）"""
        
        # 模拟异步处理
        await asyncio.sleep(0.5)
        
        base_narration = f"现在我们来看第{slide_info.page_number}页，"
        
        if slide_info.title:
            base_narration += f"这一页的主题是「{slide_info.title}」。"
        
        if slide_info.content and slide_info.content.strip():
            content_lines = [line.strip() for line in slide_info.content.split('\n') if line.strip()]
            if content_lines:
                if len(content_lines) == 1:
                    base_narration += f"主要内容是：{content_lines[0]}。"
                else:
                    base_narration += f"主要包含以下几个要点：{', '.join(content_lines[:3])}。"
        else:
            base_narration += "这一页主要是图片内容，让我们仔细观察一下图片中展示的信息。"
        
        # 根据风格调整语调
        if style == "casual":
            base_narration += "大家可以看到，这部分内容还是很有意思的。"
        elif style == "educational":
            base_narration += "这个概念对我们理解整个主题非常重要，建议大家重点关注。"
        elif style == "enthusiastic":
            base_narration += "这真是一个精彩的内容！让我们深入了解一下。"
        else:  # professional
            base_narration += "这部分内容是我们今天讨论的重点之一。"
        
        return base_narration
    
    async def get_slide_image(self, document_id: str, page_num: int) -> Optional[bytes]:
        """获取幻灯片图片数据"""
        image_filename = f"{document_id}_page_{page_num}.png"
        image_path = self.images_path / image_filename
        
        if image_path.exists():
            with open(image_path, "rb") as f:
                return f.read()
        return None
    
    async def get_slide_thumbnail(self, document_id: str, page_num: int) -> Optional[bytes]:
        """获取幻灯片缩略图数据"""
        thumbnail_filename = f"{document_id}_thumb_{page_num}.png"
        thumbnail_path = self.thumbnails_path / thumbnail_filename
        
        if thumbnail_path.exists():
            with open(thumbnail_path, "rb") as f:
                return f.read()
        return None

    async def extract_ppt_text_content(self, file_path: str) -> List[Dict[str, Any]]:
        """提取PPTX文件的完整文本内容"""
        try:
            if not PPTX_SUPPORT:
                raise Exception("python-pptx库未安装，无法提取文本内容")
            
            file_path_obj = Path(file_path)
            if not file_path_obj.exists():
                raise FileNotFoundError(f"文件不存在: {file_path}")
            
            prs = Presentation(str(file_path_obj))
            slide_contents = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # 提取标题
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                
                # 提取正文内容
                content_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # 跳过标题，避免重复
                        if shape != slide.shapes.title:
                            content_text += shape.text.strip() + "\n"
                
                # 提取演讲者备注
                notes = ""
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    for shape in notes_slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            notes += shape.text.strip() + "\n"
                
                slide_data = {
                    "page_number": slide_num,
                    "title": title or f"第{slide_num}页",
                    "content": content_text.strip(),
                    "notes": notes.strip() if notes else "",
                    "full_text": f"{title}\n{content_text}\n{notes}".strip()
                }
                
                slide_contents.append(slide_data)
                
            return slide_contents
            
        except Exception as e:
            logger.error(f"提取PPTX文本内容失败: {e}")
            raise

    async def generate_slide_narration(self, 
                                     slide_content: Dict[str, Any], 
                                     total_slides: int,
                                     narration_style: str = "讲解") -> str:
        """为单页幻灯片生成讲解文本"""
        try:
            # 初始化LLM服务
            await self._initialize_llm_if_needed()
            
            # 获取LLM配置
            available_providers = llm_manager.get_available_providers()
            if not available_providers:
                raise Exception("没有可用的LLM服务提供商")
            
            # 优先使用豆包，如果不可用则使用第一个可用的提供商
            provider = "douyin" if "douyin" in available_providers else available_providers[0]
            
            # 构建提示词
            page_num = slide_content["page_number"]
            title = slide_content["title"]
            content = slide_content["content"]
            notes = slide_content["notes"]
            
            # 组合完整文本
            full_text = ""
            if title:
                full_text += f"标题：{title}\n"
            if content:
                full_text += f"内容：{content}\n"
            if notes:
                full_text += f"备注：{notes}\n"
            
            if not full_text.strip():
                full_text = f"第{page_num}页（无文字内容）"
            
            prompt = f"""你是一位专业的PPT讲解员，现在需要为一个PPT的第{page_num}页（共{total_slides}页）生成{narration_style}文本。

PPT页面内容：
{full_text}

请根据以上内容生成一段自然流畅的{narration_style}文本，要求：
1. 文字长度严格控制在500字以内
2. 语言自然、生动，适合口语表达
3. 重点突出关键信息
4. 如果是第1页，可以适当添加开场白
5. 如果是最后一页，可以适当添加总结语
6. 避免过多的技术术语，使用通俗易懂的表达

讲解文本："""

            # 调用LLM生成
            response = llm_manager.chat(
                provider=provider,
                message=prompt,
                temperature=0.7,
                max_tokens=600  # 稍微多一点，确保500字内容完整
            )
            
            logger.info(f"LLM响应类型: {type(response)}, 内容: {response}")
            
            # 处理不同类型的响应格式
            generated_text = ""
            if isinstance(response, str):
                # 豆包等直接返回字符串的情况
                generated_text = response.strip()
            elif isinstance(response, dict):
                if not response.get('success', True):
                    raise Exception(f"LLM生成失败: {response.get('error', '未知错误')}")
                generated_text = response.get('response', '').strip()
            else:
                # 其他格式，尝试转换为字符串
                generated_text = str(response).strip()
            
            if not generated_text:
                raise Exception("LLM返回内容为空")
            
            # 确保文本长度不超过500字
            if len(generated_text) > 500:
                generated_text = generated_text[:497] + "..."
            
            return generated_text
            
        except Exception as e:
            logger.error(f"生成幻灯片讲解失败: {e}")
            # 返回默认文本
            return f"这是第{slide_content['page_number']}页的内容：{slide_content.get('title', '无标题')}。{slide_content.get('content', '')[:100]}..."

    async def synthesize_slide_audio(self, 
                                   text: str, 
                                   document_id: str, 
                                   page_number: int,
                                   voice_settings: Dict[str, Any] = None) -> Dict[str, Any]:
        """为幻灯片文本合成语音"""
        try:
            # 初始化TTS服务
            await self._initialize_tts_if_needed()
            
            # 获取TTS配置
            available_providers = tts_manager.get_available_providers()
            if not available_providers:
                raise Exception("没有可用的TTS服务提供商")
            
            # 优先使用豆包，如果不可用则使用第一个可用的提供商
            provider = "douyin" if "douyin" in available_providers else available_providers[0]
            
            # 默认语音设置
            default_settings = {
                "voice": "zh_male_beijingxiaoye_emo_v2_mars_bigtts",
                "speed": 1.0,
                "pitch": 1.0,
                "volume": 1.0,
                "format": "wav"
            }
            
            if voice_settings:
                default_settings.update(voice_settings)
            
            # 生成音频文件名
            audio_filename = f"{document_id}_page_{page_number}.{default_settings['format']}"
            audio_path = self.audio_path / audio_filename
            
            # 调用TTS合成
            response = tts_manager.synthesize_text(
                provider=provider,
                text=text,
                output_file=str(audio_path),
                voice=default_settings["voice"],
                speed=default_settings["speed"],
                pitch=default_settings["pitch"],
                volume=default_settings["volume"],
                audio_format=default_settings["format"]
            )
            
            if not response.success:
                raise Exception(f"TTS合成失败: {response.error_msg}")
            
            # 获取音频文件大小
            audio_size = 0
            if audio_path.exists():
                audio_size = audio_path.stat().st_size
            
            return {
                "success": True,
                "audio_file": audio_filename,
                "audio_path": str(audio_path),
                "audio_size": audio_size,
                "duration": response.duration if hasattr(response, 'duration') else None,
                "text": text,
                "provider": provider
            }
            
        except Exception as e:
            logger.error(f"合成语音失败: {e}")
            return {
                "success": False,
                "error": str(e),
                "text": text
            }

    async def generate_ppt_narration(self, 
                                   document_id: str,
                                   file_path: str,
                                   voice_settings: Dict[str, Any] = None,
                                   narration_style: str = "讲解") -> Dict[str, Any]:
        """为整个PPT生成讲解音频"""
        try:
            # 提取PPT文本内容
            slide_contents = await self.extract_ppt_text_content(file_path)
            total_slides = len(slide_contents)
            
            results = []
            generated_files = []
            
            for slide_data in slide_contents:
                page_num = slide_data["page_number"]
                
                logger.info(f"正在处理第{page_num}页...")
                
                # 生成讲解文本
                narration_text = await self.generate_slide_narration(
                    slide_data, 
                    total_slides, 
                    narration_style
                )
                
                # 合成语音
                audio_result = await self.synthesize_slide_audio(
                    narration_text,
                    document_id,
                    page_num,
                    voice_settings
                )
                
                page_result = {
                    "page_number": page_num,
                    "title": slide_data["title"],
                    "narration_text": narration_text,
                    "audio_result": audio_result
                }
                
                results.append(page_result)
                
                if audio_result.get("success"):
                    generated_files.append(audio_result["audio_file"])
                
                # 短暂延迟，避免API调用过快
                await asyncio.sleep(0.5)
            
            success_count = sum(1 for r in results if r["audio_result"].get("success"))
            
            return {
                "success": True,
                "document_id": document_id,
                "total_slides": total_slides,
                "success_count": success_count,
                "failed_count": total_slides - success_count,
                "results": results,
                "generated_files": generated_files,
                "narration_style": narration_style
            }
            
        except Exception as e:
            logger.error(f"生成PPT讲解失败: {e}")
            return {
                "success": False,
                "error": str(e),
                "document_id": document_id
            }

    async def get_slide_audio(self, document_id: str, page_number: int) -> Optional[bytes]:
        """获取幻灯片音频数据"""
        audio_filename = f"{document_id}_page_{page_number}.wav"
        audio_path = self.audio_path / audio_filename
        
        if audio_path.exists():
            with open(audio_path, "rb") as f:
                return f.read()
        return None

    async def _initialize_llm_if_needed(self):
        """根据需要初始化LLM服务"""
        try:
            if config_center_service:
                config = await config_center_service.get_config("llm")
                if config and config.get('enabled', True):
                    # 提取各提供商的配置
                    provider_configs = {}
                    for provider in ['aliyun', 'baidu', 'douyin', 'tencent', 'xunfei']:
                        provider_config = config.get(provider, {})
                        if provider_config and any(v for v in provider_config.values() if v):
                            provider_configs[provider] = provider_config

                    # 初始化LLM实例
                    if provider_configs:
                        llm_manager.initialize_llms(provider_configs)
                        logger.info(f"成功初始化LLM服务，可用提供商: {list(provider_configs.keys())}")
                    else:
                        logger.warning("没有找到有效的LLM配置")
        except Exception as e:
            logger.error(f"初始化LLM服务失败: {e}")

    async def _initialize_tts_if_needed(self):
        """根据需要初始化TTS服务"""
        try:
            if config_center_service:
                config = await config_center_service.get_config("tts")
                if config and config.get('enabled', True):
                    # 提取各提供商的配置
                    provider_configs = {}
                    for provider in ['baidu', 'xunfei', 'aliyun', 'tencent', 'douyin']:
                        provider_config = config.get(provider, {})
                        if provider_config and any(v for v in provider_config.values() if v):
                            provider_configs[provider] = provider_config

                    # 初始化TTS实例
                    if provider_configs:
                        tts_manager.initialize_tts_services(provider_configs)
                        logger.info(f"成功初始化TTS服务，可用提供商: {list(provider_configs.keys())}")
                    else:
                        logger.warning("没有找到有效的TTS配置")
        except Exception as e:
            logger.error(f"初始化TTS服务失败: {e}")


# 创建全局服务实例
slide_service = SlideService()
